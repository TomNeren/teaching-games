#!/usr/bin/env python3
"""
Bereinigte PowerPoint-Präsentation erstellen (ohne Vorlage-Beispiele)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

# Konfiguration
SKILL_PATH_PPTX = Path("/root/.openclaw/workspace/teaching/tom-skills-extracted/pptx/pptx-v4/pptx-v4")
OUTPUT_PATH = Path("/root/.openclaw/workspace/teaching/output/2BFH2-GGK/Holocaust")
TEMPLATE_PPTX = SKILL_PATH_PPTX / "Vorlage.pptx"
TIMER_PATH = SKILL_PATH_PPTX / "timer"

def add_timer(slide, minutes: int):
    """Fügt Timer-Video am unteren Rand ein."""
    timer_file = f"timer_pixel_{minutes}min.mp4"
    timer_path = TIMER_PATH / timer_file
    
    if not timer_path.exists():
        print(f"⚠️ Timer nicht gefunden: {timer_path}")
        return False
    
    slide.shapes.add_movie(
        str(timer_path),
        Inches(0),        # left
        Inches(13.8),     # top
        Inches(26.67),    # width
        Inches(1.2),      # height
        mime_type='video/mp4'
    )
    return True

def set_placeholder(slide, idx: int, text: str):
    """Setzt Text in Placeholder."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            shape.text = text
            return True
    return False

def fix_fonts(prs):
    """Korrigiert Keynote-Platzhalter-Fonts."""
    FONT_REPLACEMENTS = {
        '+mj-lt': 'Avenir Next',
        '+mn-lt': 'Avenir Next',
        '+mj-ea': 'Avenir Next',
        '+mn-ea': 'Avenir Next',
    }
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name in FONT_REPLACEMENTS:
                            run.font.name = FONT_REPLACEMENTS[run.font.name]
    return prs

# Neue Präsentation mit nur den Master Slides
prs = Presentation(str(TEMPLATE_PPTX))

# Alle vorhandenen Folien löschen
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

print(f"✅ Vorlage geladen, Beispiel-Folien entfernt")

# Layout-Indizes
LAYOUTS = {
    "titel": 0,
    "titel_punkte": 3,
    "arbeitsphase": 5,
    "abschnitt": 7,
    "fakt": 11,
}

# Folie 1: Titelfolie
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["titel"]])
set_placeholder(slide, 0, "Holocaust")
set_placeholder(slide, 1, "Vom Rechtsstaat zum Unrechtsstaat")

# Folie 2: Einstieg - Provokante These
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["abschnitt"]])
set_placeholder(slide, 0, '"1933 war Deutschland noch eine Demokratie."')

# Folie 3: Meinungslinie
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["titel_punkte"]])
title_shape = slide.shapes.title
title_shape.text = "Positionierung"
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "Stimmt die These?"
p = tf.add_paragraph()
p.text = "• Stimme zu"
p = tf.add_paragraph()
p.text = "• Stimme nicht zu"
p = tf.add_paragraph()
p.text = "• Bin unsicher"

# Folie 4: Lernziele
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["titel_punkte"]])
title_shape = slide.shapes.title
title_shape.text = "Lernziele"
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "Heute lernst du:"
p = tf.add_paragraph()
p.text = "• Wie aus einer Demokratie ein Unrechtsstaat wurde"
p = tf.add_paragraph()
p.text = "• Wie Propaganda im Nationalsozialismus funktionierte"
p = tf.add_paragraph()
p.text = "• Warum Zivilcourage wichtig ist"

# Folie 5: Überblick Machtergreifung
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["titel_punkte"]])
title_shape = slide.shapes.title
title_shape.text = "1933-1934: Vom Rechtsstaat zum Unrechtsstaat"
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "Schlüsselereignisse:"
p = tf.add_paragraph()
p.text = "• 30.01.1933: Hitler wird Reichskanzler"
p = tf.add_paragraph()
p.text = "• 27.02.1933: Reichstagsbrand → Notverordnungen"
p = tf.add_paragraph()
p.text = "• 23.03.1933: Ermächtigungsgesetz"
p = tf.add_paragraph()
p.text = "• 1933-34: Gleichschaltung (Parteien, Gewerkschaften, Medien)"
p = tf.add_paragraph()
p.text = "• 02.08.1934: Hitler wird 'Führer und Reichskanzler'"

# Folie 6: Key Message
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["fakt"]])
for shape in slide.shapes:
    if shape.has_text_frame:
        shape.text = "In nur 18 Monaten wurde Deutschland von einer Demokratie zur Diktatur."

# Folie 7: Arbeitsphase - Propaganda (mit 5min Timer)
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["arbeitsphase"]])
title_shape = slide.shapes.title
title_shape.text = "Arbeitsauftrag: Propaganda-Analyse"
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "Analysiere das NS-Propagandaplakat:"
p = tf.add_paragraph()
p.text = "1. Wer ist die Zielgruppe?"
p = tf.add_paragraph()
p.text = "2. Welche Emotionen werden angesprochen?"
p = tf.add_paragraph()
p.text = "3. Welche Feindbilder werden gezeigt?"
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "→ Arbeite mit deinem Partner (5 Minuten)"
add_timer(slide, 5)

# Folie 8: Gruppenpuzzle - Eskalationsstufen
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["arbeitsphase"]])
title_shape = slide.shapes.title
title_shape.text = "Gruppenpuzzle: Vom Ausschluss zur Vernichtung"
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "4 Stationen (Expertengruppen):"
p = tf.add_paragraph()
p.text = "1. Nürnberger Gesetze (1935)"
p = tf.add_paragraph()
p.text = "2. Reichspogromnacht (1938)"
p = tf.add_paragraph()
p.text = "3. Ghettos (1940-1942)"
p = tf.add_paragraph()
p.text = "4. Vernichtungslager (1942-1945)"
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "→ Lies dein Material (4 Minuten)"
add_timer(slide, 4)

# Folie 9: Stammgruppen (mit 10min Timer)
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["arbeitsphase"]])
title_shape = slide.shapes.title
title_shape.text = "Stammgruppen: Austausch"
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "Teilt euer Wissen:"
p = tf.add_paragraph()
p.text = "• Jede*r erklärt seine Station (2 Min)"
p = tf.add_paragraph()
p.text = "• Was war die Eskalation?"
p = tf.add_paragraph()
p.text = "• Wer hätte etwas tun können?"
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "→ Austausch in der Gruppe (10 Minuten)"
add_timer(slide, 10)

# Folie 10: Reflexion - Gegenwartsbezug
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["abschnitt"]])
set_placeholder(slide, 0, "Und heute?")

# Folie 11: Blitzlicht
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["titel_punkte"]])
title_shape = slide.shapes.title
title_shape.text = "Blitzlicht: Mechanismen heute"
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "Welche Mechanismen erkenne ich heute?"
p = tf.add_paragraph()
p.text = "• Autoritäre Tendenzen?"
p = tf.add_paragraph()
p.text = "• Propaganda & Fake News?"
p = tf.add_paragraph()
p.text = "• Ausgrenzung & Diskriminierung?"
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "→ Nenne 1 Beispiel"

# Folie 12: Schluss
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["fakt"]])
for shape in slide.shapes:
    if shape.has_text_frame:
        shape.text = "Demokratie ist nicht selbstverständlich. Sie muss jeden Tag verteidigt werden."

# Font fix
prs = fix_fonts(prs)

# Speichern
output_file = OUTPUT_PATH / "Stunde_01_PPT_clean.pptx"
prs.save(str(output_file))
print(f"✅ Bereinigte PowerPoint erstellt: {output_file}")
print(f"   Anzahl Folien: {len(prs.slides)}")
