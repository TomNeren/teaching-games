#!/usr/bin/env python3
"""
Holocaust-Unterrichtsstunde für 2BFH2-GGK erstellen
Erstellt: Stundenplanung (MD), PowerPoint mit Timer, 4 Arbeitsblätter
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document
from pathlib import Path
import os
import re

# ==================== KONFIGURATION ====================

# Pfade
SKILL_PATH_PPTX = Path("/root/.openclaw/workspace/teaching/tom-skills-extracted/pptx/pptx-v4/pptx-v4")
SKILL_PATH_DOC = Path("/root/.openclaw/workspace/teaching/tom-skills-extracted/doc/arbeitsblatt-erstellen-v3/arbeitsblatt-erstellen-v3")
OUTPUT_PATH = Path("/root/.openclaw/workspace/teaching/output/2BFH2-GGK/Holocaust")

TEMPLATE_PPTX = SKILL_PATH_PPTX / "Vorlage.pptx"
TEMPLATE_DOC = SKILL_PATH_DOC / "templates" / "Vorlage_Fach.docx"
TIMER_PATH = SKILL_PATH_PPTX / "timer"

# Ausgabepfade
OUTPUT_PATH.mkdir(parents=True, exist_ok=True)
(OUTPUT_PATH / "doc").mkdir(exist_ok=True)

# Farbpalette: Midnight Executive (ernst, dunkel)
COLOR_PRIMARY = RGBColor(0x1E, 0x27, 0x61)
COLOR_SECONDARY = RGBColor(0xCA, 0xDC, 0xFC)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ==================== HILFSFUNKTIONEN ====================

def replace_header_placeholders(doc, replacements: dict):
    """Ersetzt Platzhalter in Header-Tabellen."""
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    full_text = "".join(run.text for run in paragraph.runs)
                    for old, new in replacements.items():
                        if old in full_text:
                            full_text = full_text.replace(old, new)
                    if paragraph.runs:
                        paragraph.runs[0].text = full_text
                        for run in paragraph.runs[1:]:
                            run.text = ""
    return doc

def clear_body(doc):
    """Entfernt alle Paragraphen nach dem ersten."""
    while len(doc.paragraphs) > 1:
        p = doc.paragraphs[-1]._element
        p.getparent().remove(p)
    if doc.paragraphs:
        doc.paragraphs[0].clear()
    return doc

def validate_document(doc_path: str) -> bool:
    """Validiert ein erstelltes Dokument."""
    doc = Document(doc_path)
    errors = []

    # Gesamttext sammeln
    all_text = ""
    for p in doc.paragraphs:
        all_text += p.text + "\n"
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                all_text += cell.text + "\n"

    # Check 1: Nicht ersetzte Platzhalter
    placeholders = re.findall(r'\[[^\]]+\]', all_text)
    if placeholders:
        errors.append(f"❌ PLATZHALTER NICHT ERSETZT: {placeholders}")

    # Check 2: Falsche Umlaute
    umlaut_errors = ['ae', 'oe', 'ue', 'Ae', 'Oe', 'Ue']
    for ue in umlaut_errors:
        if ue in all_text:
            errors.append(f"❌ UMLAUT-FEHLER: '{ue}' gefunden")

    # Ergebnis
    if errors:
        print("\n".join(errors))
        return False
    print(f"✅ Validiert: {doc_path}")
    return True

def add_timer(slide, minutes: int):
    """Fügt Timer-Video am unteren Rand ein."""
    timer_file = f"timer_pixel_{minutes}min.mp4"
    timer_path = TIMER_PATH / timer_file
    
    if not timer_path.exists():
        print(f"⚠️ Timer nicht gefunden: {timer_path}")
        return False
    
    slide.shapes.add_movie(
        str(timer_path),
        Inches(0),        # left
        Inches(13.8),     # top
        Inches(26.67),    # width
        Inches(1.2),      # height
        mime_type='video/mp4'
    )
    return True

def set_placeholder(slide, idx: int, text: str):
    """Setzt Text in Placeholder."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            shape.text = text
            return True
    return False

def fix_fonts(prs):
    """Korrigiert Keynote-Platzhalter-Fonts."""
    FONT_REPLACEMENTS = {
        '+mj-lt': 'Avenir Next',
        '+mn-lt': 'Avenir Next',
        '+mj-ea': 'Avenir Next',
        '+mn-ea': 'Avenir Next',
    }
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name in FONT_REPLACEMENTS:
                            run.font.name = FONT_REPLACEMENTS[run.font.name]
    return prs

# ==================== STUNDENPLANUNG (MARKDOWN) ====================

def create_stundenplanung():
    """Erstellt die Stundenplanung als Markdown."""
    content = """# Stundenplanung: Holocaust – Vom Rechtsstaat zum Unrechtsstaat

**Klasse**: 2BFH2  
**Fach**: Geschichte mit Gemeinschaftskunde (GGK)  
**Dauer**: 90 Minuten (Doppelstunde)  
**Thema/Lernziel**: Holocaust – Mechanismen von Verführung, Propaganda und Gewalt verstehen  
**Kompetenzen**: BPE 2.1 "Wie lässt sich Demokratie schützen und erneuern?" — Nationalsozialistische Diktatur, Propaganda, Zivilcourage  
**Leitfrage**: "Wie konnte aus einer Demokratie ein Unrechtsstaat werden — und was hat das mit uns zu tun?"

---

## Materialien

- PowerPoint-Präsentation mit Timer
- AB 01: Propaganda-Analyse (Niveau A + B)
- AB 02: Perspektivwechsel – Zivilcourage (Niveau A + B)
- iPads für Recherche
- Beamer

---

## Stundenverlauf

| Phase | Zeit | Methode | Beschreibung | Material |
|-------|------|---------|--------------|----------|
| **Einstieg** | 10 min | Bildimpuls + These | Provokante These: "1933 war Deutschland noch eine Demokratie." SuS positionieren sich (Meinungslinie). Kurze Diskussion: Was ist Demokratie? Wie kann sie enden? | PPT Folie 2-3 |
| **Lernziele** | 5 min | Überblick | Lernziele transparent machen + Leitfrage einführen | PPT Folie 4 |
| **Erarbeitung I** | 15 min | Lehrer-Input + Zeitstrahl | Überblick: Machtergreifung → Gleichschaltung (1933-1934). Visualisierung als Zeitstrahl. SuS notieren 3-5 Schlüsselereignisse. | PPT Folie 5-6, Tafel |
| **Erarbeitung II** | 25 min | Think-Pair-Share + Quellenanalyse | **Hauptteil**: Propaganda als Werkzeug. SuS analysieren NS-Propagandaplakate in PA (5min Timer). Leitfragen: Wer ist die Zielgruppe? Welche Emotionen? Welche Feindbilder? Austausch im Plenum. | AB 01 (Niveau A/B), PPT Folie 7 (mit 5min Timer) |
| **Erarbeitung III** | 20 min | Gruppenpuzzle (Jigsaw) | Vom Ausschluss zur Vernichtung: 4 Stationen (Nürnberger Gesetze, Reichspogromnacht, Ghettos, Vernichtungslager). Expertengruppen (5min), dann Stammgruppen (10min). | PPT Folie 8-9 (mit 4min Timer), Infomaterial (digital) |
| **Sicherung** | 10 min | Schreibgespräch | Brief aus Perspektive eines Zeitzeugen schreiben (3 Sätze): "Was ich erlebt habe..." SuS schreiben anonym auf Karteikarten, werden vorgelesen. | AB 02 (Niveau A/B) |
| **Reflexion** | 5 min | Blitzlicht | Gegenwartsbezug: "Welche Mechanismen erkenne ich heute?" — Autoritäre Tendenzen, Propaganda, Fake News. SuS nennen je 1 Beispiel. | PPT Folie 10 |

**Gesamtzeit**: 90 Minuten

---

## Differenzierung

### Niveau Niedrig (~40% der SuS)
- **AB 01 (Niveau A)**: Propaganda-Analyse mit Wortbanken, vorgegebenen Satzanfängen, Scaffolding
- **AB 02 (Niveau A)**: Lückentext + geführtes Schreiben mit Satzbausteinen
- **Erarbeitung III**: Vereinfachte Infotexte, Bilder als Unterstützung
- **Sicherung**: Satzbausteine für Brief ("Ich habe erlebt, dass...")

### Niveau Mittel (~40% der SuS)
- **AB 01 (Niveau B)**: Propaganda-Analyse mit offenen Leitfragen
- **AB 02 (Niveau B)**: Freier Brief aus Zeitzeugen-Perspektive
- **Erarbeitung III**: Standard-Infotexte

### Niveau Hoch (~20% der SuS)
- **Zusatzaufgabe**: Vergleich zu heutigen Propaganda-Strategien (Social Media)
- **Transfer**: "Was können wir heute tun, um Demokratie zu schützen?"
- **Reflexion**: Weiterführende Diskussion zu Zivilcourage

---

## Gegenwartsbezug

- **Autoritäre Tendenzen weltweit**: Vergleich zu heutigen Entwicklungen (Türkei, Russland, Ungarn)
- **Fake News & Social Media**: Wie funktioniert moderne Propaganda?
- **Zivilcourage heute**: Beispiele aus Gegenwart (z.B. Sea Watch, Fridays for Future, Proteste in Iran)

---

## Leistungserhebung

- Mündliche Beteiligung (Diskussion, Präsentation)
- Schriftliche Quellenanalyse (AB 01)
- Brief als kreative Leistung (AB 02)

---

## Didaktische Hinweise

- **Yad Vashem Guidelines beachten**: Opfer als Menschen zeigen, nicht nur als Statistik
- **Sensibilität**: Kulturelle/religiöse Diversität in der Klasse berücksichtigen
- **Keine Überwältigung**: Holocaust sachlich, aber empathisch behandeln
- **Handlungsorientierung**: SuS zu Zivilcourage ermutigen, nicht nur Wissen vermitteln

---

*Erstellt: 2026-02-03*  
*Skill-Version: unterrichtsplanung-workflow v3 + pptx v4 + arbeitsblatt v3*
"""
    
    output_file = OUTPUT_PATH / "Stundenplanung_Holocaust.md"
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(content)
    print(f"✅ Stundenplanung erstellt: {output_file}")
    return str(output_file)

# ==================== POWERPOINT-PRÄSENTATION ====================

def create_powerpoint():
    """Erstellt PowerPoint-Präsentation mit Timer."""
    prs = Presentation(str(TEMPLATE_PPTX))
    
    # Layout-Indizes (aus unterricht.md)
    LAYOUTS = {
        "titel": 0,
        "titel_punkte": 3,
        "titel_punkte_foto": 4,
        "arbeitsphase": 5,
        "abschnitt": 7,
        "fakt": 11,
    }
    
    # Folie 1: Titelfolie
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["titel"]])
    set_placeholder(slide, 0, "Holocaust")
    set_placeholder(slide, 1, "Vom Rechtsstaat zum Unrechtsstaat")
    
    # Folie 2: Einstieg - Provokante These
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["abschnitt"]])
    set_placeholder(slide, 0, '"1933 war Deutschland noch eine Demokratie."')
    
    # Folie 3: Meinungslinie
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["titel_punkte"]])
    title_shape = slide.shapes.title
    title_shape.text = "Positionierung"
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Stimmt die These?"
    p = tf.add_paragraph()
    p.text = "• Stimme zu"
    p = tf.add_paragraph()
    p.text = "• Stimme nicht zu"
    p = tf.add_paragraph()
    p.text = "• Bin unsicher"
    
    # Folie 4: Lernziele
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["titel_punkte"]])
    title_shape = slide.shapes.title
    title_shape.text = "Lernziele"
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Heute lernst du:"
    p = tf.add_paragraph()
    p.text = "• Wie aus einer Demokratie ein Unrechtsstaat wurde"
    p = tf.add_paragraph()
    p.text = "• Wie Propaganda im Nationalsozialismus funktionierte"
    p = tf.add_paragraph()
    p.text = "• Warum Zivilcourage wichtig ist"
    
    # Folie 5: Überblick Machtergreifung
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["titel_punkte"]])
    title_shape = slide.shapes.title
    title_shape.text = "1933-1934: Vom Rechtsstaat zum Unrechtsstaat"
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Schlüsselereignisse:"
    p = tf.add_paragraph()
    p.text = "• 30.01.1933: Hitler wird Reichskanzler"
    p = tf.add_paragraph()
    p.text = "• 27.02.1933: Reichstagsbrand → Notverordnungen"
    p = tf.add_paragraph()
    p.text = "• 23.03.1933: Ermächtigungsgesetz"
    p = tf.add_paragraph()
    p.text = "• 1933-34: Gleichschaltung (Parteien, Gewerkschaften, Medien)"
    p = tf.add_paragraph()
    p.text = "• 02.08.1934: Hitler wird 'Führer und Reichskanzler'"
    
    # Folie 6: Key Message
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["fakt"]])
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = "In nur 18 Monaten wurde Deutschland von einer Demokratie zur Diktatur."
    
    # Folie 7: Arbeitsphase - Propaganda (mit 5min Timer)
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["arbeitsphase"]])
    title_shape = slide.shapes.title
    title_shape.text = "Arbeitsauftrag: Propaganda-Analyse"
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Analysiere das NS-Propagandaplakat:"
    p = tf.add_paragraph()
    p.text = "1. Wer ist die Zielgruppe?"
    p = tf.add_paragraph()
    p.text = "2. Welche Emotionen werden angesprochen?"
    p = tf.add_paragraph()
    p.text = "3. Welche Feindbilder werden gezeigt?"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "→ Arbeite mit deinem Partner (5 Minuten)"
    add_timer(slide, 5)
    
    # Folie 8: Gruppenpuzzle - Eskalationsstufen
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["arbeitsphase"]])
    title_shape = slide.shapes.title
    title_shape.text = "Gruppenpuzzle: Vom Ausschluss zur Vernichtung"
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "4 Stationen (Expertengruppen):"
    p = tf.add_paragraph()
    p.text = "1. Nürnberger Gesetze (1935)"
    p = tf.add_paragraph()
    p.text = "2. Reichspogromnacht (1938)"
    p = tf.add_paragraph()
    p.text = "3. Ghettos (1940-1942)"
    p = tf.add_paragraph()
    p.text = "4. Vernichtungslager (1942-1945)"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "→ Lies dein Material (4 Minuten)"
    add_timer(slide, 4)
    
    # Folie 9: Stammgruppen (mit 10min Timer)
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["arbeitsphase"]])
    title_shape = slide.shapes.title
    title_shape.text = "Stammgruppen: Austausch"
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Teilt euer Wissen:"
    p = tf.add_paragraph()
    p.text = "• Jede*r erklärt seine Station (2 Min)"
    p = tf.add_paragraph()
    p.text = "• Was war die Eskalation?"
    p = tf.add_paragraph()
    p.text = "• Wer hätte etwas tun können?"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "→ Austausch in der Gruppe (10 Minuten)"
    add_timer(slide, 10)
    
    # Folie 10: Reflexion - Gegenwartsbezug
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["abschnitt"]])
    set_placeholder(slide, 0, "Und heute?")
    
    # Folie 11: Blitzlicht
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["titel_punkte"]])
    title_shape = slide.shapes.title
    title_shape.text = "Blitzlicht: Mechanismen heute"
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Welche Mechanismen erkenne ich heute?"
    p = tf.add_paragraph()
    p.text = "• Autoritäre Tendenzen?"
    p = tf.add_paragraph()
    p.text = "• Propaganda & Fake News?"
    p = tf.add_paragraph()
    p.text = "• Ausgrenzung & Diskriminierung?"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "→ Nenne 1 Beispiel"
    
    # Folie 12: Schluss
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["fakt"]])
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = "Demokratie ist nicht selbstverständlich. Sie muss jeden Tag verteidigt werden."
    
    # Font fix
    prs = fix_fonts(prs)
    
    # Speichern
    output_file = OUTPUT_PATH / "Stunde_01_PPT.pptx"
    prs.save(str(output_file))
    print(f"✅ PowerPoint erstellt: {output_file}")
    return str(output_file)

# ==================== ARBEITSBLÄTTER ====================

def create_arbeitsblatt_01_niveau_b():
    """AB 01: Propaganda-Analyse - Niveau B (Standard)"""
    doc = Document(str(TEMPLATE_DOC))
    
    replacements = {
        '[Thema]': 'NS-Propaganda: Wie funktioniert Verführung?',
        '[Fach]': 'GGK',
        '[Ziel 1 – nicht länger als eine Zeile]': 'Ich kann Propagandaplakate analysieren',
        '[Ziel 2 – nicht länger als eine Zeile]': 'Ich erkenne Manipulationstechniken',
        '[Ziel 3 – nicht länger als eine Zeile]': 'Ich vergleiche mit heutigen Strategien',
        '[A / B / C]': 'B',
    }
    doc = replace_header_placeholders(doc, replacements)
    doc = clear_body(doc)
    
    # Inhalt
    p = doc.add_paragraph("Teil 1: Quellenanalyse", style='Heading 1')
    p.runs[0].underline = True
    
    doc.add_paragraph("""Analysiere das NS-Propagandaplakat (wird im Unterricht gezeigt) und beantworte die folgenden Fragen:""")
    
    doc.add_paragraph("")
    doc.add_paragraph("1. Beschreibung:")
    doc.add_paragraph("   Was siehst du auf dem Plakat? (Personen, Symbole, Farben)")
    doc.add_paragraph("")
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("")
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("")
    
    doc.add_paragraph("2. Zielgruppe:")
    doc.add_paragraph("   Wen soll das Plakat ansprechen?")
    doc.add_paragraph("")
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("")
    
    doc.add_paragraph("3. Emotionen:")
    doc.add_paragraph("   Welche Gefühle sollen beim Betrachter ausgelöst werden?")
    doc.add_paragraph("")
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("")
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("")
    
    doc.add_paragraph("4. Feindbilder:")
    doc.add_paragraph("   Wer wird als Feind dargestellt?")
    doc.add_paragraph("")
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("")
    
    p = doc.add_paragraph("Teil 2: Transfer", style='Heading 1')
    p.runs[0].underline = True
    
    doc.add_paragraph("5. Vergleich:")
    doc.add_paragraph("   Kennst du ähnliche Strategien aus heutigen Medien? (Social Media, Werbung, Politik)")
    doc.add_paragraph("")
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("")
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("")
    
    output_file = OUTPUT_PATH / "doc" / "AB_01_Propaganda_NiveauB.docx"
    doc.save(str(output_file))
    validate_document(str(output_file))
    return str(output_file)

def create_arbeitsblatt_01_niveau_a():
    """AB 01: Propaganda-Analyse - Niveau A (Vereinfacht)"""
    doc = Document(str(TEMPLATE_DOC))
    
    replacements = {
        '[Thema]': 'NS-Propaganda: Wie funktioniert Verführung?',
        '[Fach]': 'GGK',
        '[Ziel 1 – nicht länger als eine Zeile]': 'Ich kann Propagandaplakate analysieren',
        '[Ziel 2 – nicht länger als eine Zeile]': 'Ich erkenne Manipulationstechniken',
        '[Ziel 3 – nicht länger als eine Zeile]': 'Ich vergleiche mit heutigen Strategien',
        '[A / B / C]': 'A',
    }
    doc = replace_header_placeholders(doc, replacements)
    doc = clear_body(doc)
    
    # Inhalt mit Scaffolding
    p = doc.add_paragraph("Teil 1: Quellenanalyse (mit Hilfen)", style='Heading 1')
    p.runs[0].underline = True
    
    doc.add_paragraph("""Analysiere das NS-Propagandaplakat und beantworte die Fragen. Nutze die Satzanfänge!""")
    
    doc.add_paragraph("")
    doc.add_paragraph("1. Beschreibung:")
    doc.add_paragraph("   Auf dem Plakat sehe ich ________________________________")
    doc.add_paragraph("")
    doc.add_paragraph("   Die Farben sind ________________________________")
    doc.add_paragraph("")
    
    doc.add_paragraph("2. Zielgruppe:")
    doc.add_paragraph("   Das Plakat richtet sich an ________________________________")
    doc.add_paragraph("")
    doc.add_paragraph("   (z.B. junge Männer, Familien, Arbeiter, ...)")
    doc.add_paragraph("")
    
    doc.add_paragraph("3. Emotionen:")
    doc.add_paragraph("   Das Plakat will folgende Gefühle auslösen:")
    doc.add_paragraph("")
    doc.add_paragraph("   ☐ Stolz        ☐ Angst        ☐ Hoffnung")
    doc.add_paragraph("   ☐ Wut          ☐ Vertrauen    ☐ Zusammengehörigkeit")
    doc.add_paragraph("")
    
    doc.add_paragraph("4. Feindbilder:")
    doc.add_paragraph("   Als Feind wird dargestellt: ________________________________")
    doc.add_paragraph("")
    
    p = doc.add_paragraph("Teil 2: Wortbank", style='Heading 1')
    p.runs[0].underline = True
    
    doc.add_paragraph("""Hilfreiche Wörter:
• Gemeinschaft • Heldentum • Bedrohung • Ausgrenzung • Stärke • Einheit
• Opfer • Täter • Manipulation • Emotion • Symbol • Uniform""")
    
    doc.add_paragraph("")
    
    p = doc.add_paragraph("Teil 3: Transfer", style='Heading 1')
    p.runs[0].underline = True
    
    doc.add_paragraph("5. Heute:")
    doc.add_paragraph("   Ähnliche Strategien sehe ich bei ________________________________")
    doc.add_paragraph("")
    doc.add_paragraph("   (z.B. auf Instagram, in der Werbung, in Nachrichtensendungen)")
    doc.add_paragraph("")
    
    output_file = OUTPUT_PATH / "doc" / "AB_01_Propaganda_NiveauA.docx"
    doc.save(str(output_file))
    validate_document(str(output_file))
    return str(output_file)

def create_arbeitsblatt_02_niveau_b():
    """AB 02: Perspektivwechsel - Niveau B (Standard)"""
    doc = Document(str(TEMPLATE_DOC))
    
    replacements = {
        '[Thema]': 'Perspektivwechsel: Zivilcourage im Nationalsozialismus',
        '[Fach]': 'GGK',
        '[Ziel 1 – nicht länger als eine Zeile]': 'Ich kann mich in Zeitzeugen hineinversetzen',
        '[Ziel 2 – nicht länger als eine Zeile]': 'Ich reflektiere über Zivilcourage',
        '[Ziel 3 – nicht länger als eine Zeile]': 'Ich übertrage es auf heute',
        '[A / B / C]': 'B',
    }
    doc = replace_header_placeholders(doc, replacements)
    doc = clear_body(doc)
    
    # Inhalt
    p = doc.add_paragraph("Aufgabe: Brief aus der Vergangenheit", style='Heading 1')
    p.runs[0].underline = True
    
    doc.add_paragraph("""Stell dir vor, du lebst im Jahr 1938 in Deutschland. Du bist Zeuge der Reichspogromnacht geworden. Schreibe einen Brief an eine Person deines Vertrauens (Freund*in, Verwandte*r im Ausland), in dem du beschreibst:""")
    
    doc.add_paragraph("")
    doc.add_paragraph("• Was du erlebt hast")
    doc.add_paragraph("• Wie du dich gefühlt hast")
    doc.add_paragraph("• Warum du nichts tun konntest (oder doch etwas getan hast)")
    doc.add_paragraph("")
    
    doc.add_paragraph("Schreibe mindestens 150 Wörter. Nutze dein Wissen aus dem Gruppenpuzzle.")
    doc.add_paragraph("")
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("")
    
    p = doc.add_paragraph("Reflexion", style='Heading 1')
    p.runs[0].underline = True
    
    doc.add_paragraph("Was können wir heute tun, um Demokratie zu schützen?")
    doc.add_paragraph("")
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("")
    
    output_file = OUTPUT_PATH / "doc" / "AB_02_Perspektivwechsel_NiveauB.docx"
    doc.save(str(output_file))
    validate_document(str(output_file))
    return str(output_file)

def create_arbeitsblatt_02_niveau_a():
    """AB 02: Perspektivwechsel - Niveau A (Vereinfacht)"""
    doc = Document(str(TEMPLATE_DOC))
    
    replacements = {
        '[Thema]': 'Perspektivwechsel: Zivilcourage im Nationalsozialismus',
        '[Fach]': 'GGK',
        '[Ziel 1 – nicht länger als eine Zeile]': 'Ich kann mich in Zeitzeugen hineinversetzen',
        '[Ziel 2 – nicht länger als eine Zeile]': 'Ich reflektiere über Zivilcourage',
        '[Ziel 3 – nicht länger als eine Zeile]': 'Ich übertrage es auf heute',
        '[A / B / C]': 'A',
    }
    doc = replace_header_placeholders(doc, replacements)
    doc = clear_body(doc)
    
    # Inhalt mit Lückentext
    p = doc.add_paragraph("Aufgabe: Brief aus der Vergangenheit (mit Hilfen)", style='Heading 1')
    p.runs[0].underline = True
    
    doc.add_paragraph("""Ergänze den Brief. Nutze dein Wissen aus dem Gruppenpuzzle.""")
    
    doc.add_paragraph("")
    doc.add_paragraph("Liebe(r) ____________________,")
    doc.add_paragraph("")
    doc.add_paragraph("ich schreibe dir, weil ich etwas Schreckliches erlebt habe. In der Nacht")
    doc.add_paragraph("")
    doc.add_paragraph("vom 9. auf den 10. November 1938 habe ich ____________________")
    doc.add_paragraph("")
    doc.add_paragraph("gesehen. Überall wurden ____________________")
    doc.add_paragraph("")
    doc.add_paragraph("zerstört. Die Menschen haben ____________________.")
    doc.add_paragraph("")
    doc.add_paragraph("")
    doc.add_paragraph("Ich habe mich ____________________")
    doc.add_paragraph("")
    doc.add_paragraph("gefühlt, weil ____________________.")
    doc.add_paragraph("")
    doc.add_paragraph("")
    doc.add_paragraph("Ich konnte nichts tun, weil ____________________")
    doc.add_paragraph("")
    doc.add_paragraph("____________________.")
    doc.add_paragraph("")
    doc.add_paragraph("Ich hoffe, dass ____________________.")
    doc.add_paragraph("")
    doc.add_paragraph("")
    doc.add_paragraph("Dein(e) ____________________")
    doc.add_paragraph("")
    
    p = doc.add_paragraph("Wortbank", style='Heading 1')
    p.runs[0].underline = True
    
    doc.add_paragraph("""Hilfreiche Wörter:
• Synagogen • Geschäfte • jüdische Familien • Angst • Scham • Wut
• verhaftet • geschlagen • weggeschaut • Polizei • SA • Nachbarn
• Angst vor Strafe • niemand geholfen hat • es sich bessert""")
    
    doc.add_paragraph("")
    
    p = doc.add_paragraph("Reflexion", style='Heading 1')
    p.runs[0].underline = True
    
    doc.add_paragraph("Was können wir heute tun?")
    doc.add_paragraph("")
    doc.add_paragraph("Wir können ____________________ tun,")
    doc.add_paragraph("")
    doc.add_paragraph("damit ____________________ nie wieder passiert.")
    doc.add_paragraph("")
    
    output_file = OUTPUT_PATH / "doc" / "AB_02_Perspektivwechsel_NiveauA.docx"
    doc.save(str(output_file))
    validate_document(str(output_file))
    return str(output_file)

# ==================== HAUPTPROGRAMM ====================

def main():
    """Erstellt alle Materialien."""
    print("=" * 80)
    print("HOLOCAUST-UNTERRICHTSSTUNDE FÜR 2BFH2-GGK")
    print("=" * 80)
    print()
    
    created_files = []
    
    # 1. Stundenplanung
    print("📝 Erstelle Stundenplanung...")
    created_files.append(create_stundenplanung())
    print()
    
    # 2. PowerPoint
    print("🎨 Erstelle PowerPoint-Präsentation...")
    created_files.append(create_powerpoint())
    print()
    
    # 3. Arbeitsblätter
    print("📄 Erstelle Arbeitsblätter...")
    created_files.append(create_arbeitsblatt_01_niveau_b())
    created_files.append(create_arbeitsblatt_01_niveau_a())
    created_files.append(create_arbeitsblatt_02_niveau_b())
    created_files.append(create_arbeitsblatt_02_niveau_a())
    print()
    
    # Zusammenfassung
    print("=" * 80)
    print("✅ ALLE MATERIALIEN ERSTELLT")
    print("=" * 80)
    print("\nErstellte Dateien:")
    for i, file in enumerate(created_files, 1):
        print(f"{i}. {file}")
    print()
    print(f"📁 Ausgabeordner: {OUTPUT_PATH}")
    print()

if __name__ == "__main__":
    main()
