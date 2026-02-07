#!/usr/bin/env python3
"""
PowerPoint-Erstellung für Holocaust-Stunde
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Erstelle Präsentation
prs = Presentation()
prs.slide_width = Inches(10)  # 16:9 Format
prs.slide_height = Inches(5.625)

# Dunkles Farbschema
DARK_BG = RGBColor(30, 30, 30)
LIGHT_TEXT = RGBColor(240, 240, 240)
ACCENT_COLOR = RGBColor(200, 50, 50)  # Dunkles Rot für Ernst

def add_title_slide(prs, title, subtitle):
    """Fügt Titelfolie hinzu"""
    slide_layout = prs.slide_layouts[6]  # Leeres Layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Hintergrund dunkel
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(9), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = LIGHT_TEXT
    title_para.alignment = PP_ALIGN.CENTER
    
    # Untertitel
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.2), Inches(9), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(180, 180, 180)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content, is_dark=True):
    """Fügt Inhaltsfolie hinzu"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Hintergrund
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG if is_dark else RGBColor(245, 245, 245)
    
    text_color = LIGHT_TEXT if is_dark else RGBColor(30, 30, 30)
    
    # Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = text_color
    
    # Inhalt
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.8)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.text = content
    
    for para in content_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = text_color
        para.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Fügt Zwei-Spalten-Folie hinzu"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Hintergrund
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = LIGHT_TEXT
    
    # Linke Spalte (David)
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(4.2), Inches(3.8)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    left_frame.text = left_content
    for para in left_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(100, 150, 255)  # Blau für David
        para.space_after = Pt(10)
    
    # Rechte Spalte (Werner)
    right_box = slide.shapes.add_textbox(
        Inches(5.3), Inches(1.3), Inches(4.2), Inches(3.8)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    right_frame.text = right_content
    for para in right_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(150, 100, 50)  # Braun für Werner
        para.space_after = Pt(10)
    
    # Trennlinie
    slide.shapes.add_shape(
        1,  # Line
        Inches(5), Inches(1.3), Inches(0.05), Inches(3.8)
    ).fill.solid()
    
    return slide

# Folie 1: Titelfolie
add_title_slide(prs, "Entscheidungen 1933-1945", "Wie konnten normale Menschen zuschauen?")

# Folie 2: Das Foto
add_content_slide(prs, "Das Foto", 
    "Boykott jüdischer Geschäfte\n1. April 1933\n\n"
    "Passanten gehen vorbei...\n\n"
    "[Hier historisches Foto einfügen]")

# Folie 3: Stilles Schreiben
add_content_slide(prs, "Stilles Schreiben (3 Min.)",
    "• Was siehst du?\n\n"
    "• Was denkst du?\n\n"
    "• Was fühlst du?\n\n\n"
    "→ Notiere deine Gedanken")

# Folie 4: Leitfrage
add_content_slide(prs, "Unsere Leitfrage",
    '"Wie konnte es dazu kommen,\n'
    'dass normale Menschen\n'
    'zuschauten?"\n\n\n'
    "Heute erkunden wir das durch zwei Perspektiven...")

# Folie 5: Perspektiven wählen
add_two_column_slide(prs, "Wähle deine Perspektive",
    "🔵 David\n\n"
    "Jüdischer Junge\n"
    "Opferperspektive\n\n"
    "Erlebt:\n"
    "• Ausgrenzung\n"
    "• Entrechtung\n"
    "• Angst",
    
    "🟤 Werner\n\n"
    "Deutscher Hitlerjunge\n"
    "Zuschauerperspektive\n\n"
    "Erlebt:\n"
    "• Gruppendruck\n"
    "• Konformität\n"
    "• Gleichgültigkeit?")

# Folie 6: Self-Explanation Prompts
add_content_slide(prs, "Während du spielst...",
    "Beantworte auf dem Arbeitsblatt:\n\n"
    "• Warum hast du dich so entschieden?\n\n"
    "• Was hätte passieren können, wenn du anders gewählt hättest?\n\n"
    "• Was hat dein Charakter gefühlt?")

# Folie 7: Perspektivwechsel
add_content_slide(prs, "Perspektivwechsel",
    "Finde einen Partner, der die ANDERE Perspektive gespielt hat.\n\n"
    "Erzählt euch gegenseitig:\n"
    "• Was ist dir passiert?\n"
    "• Welche Entscheidungen musstest du treffen?\n"
    "• Wie hat sich dein Charakter gefühlt?\n\n"
    "→ 8 Minuten (4 Min. pro Person)")

# Folie 8: Reflexion
add_content_slide(prs, 'Reflexion: "Und ich?"',
    'Stilles Schreiben (10 Minuten):\n\n'
    '"Warum ist es so schwer,\n'
    'sich gegen die Mehrheit zu stellen?"\n\n\n'
    "Optional:\n"
    '"Wo sehe ich heute ähnliche Mechanismen?"')

# Folie 9: Exit-Ticket
add_content_slide(prs, "Exit-Ticket",
    "Schreibe auf deine Karteikarte:\n\n\n"
    '"Das Wichtigste,\n'
    'das ich heute gelernt habe..."')

# Folie 10: Ausblick
add_content_slide(prs, "Ausblick: Nächste Stunde",
    "Von der Ausgrenzung zur Vernichtung\n\n"
    "Wie wurde aus Zuschauen\n"
    "Mittäterschaft?")

# Speichern
output_path = "/root/.openclaw/workspace/teaching/output/2BFH2-GGK/Holocaust_v2/Stunde_01_PPT.pptx"
prs.save(output_path)
print(f"✓ Präsentation erstellt: {output_path}")
